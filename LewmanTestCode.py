# Cheryl Lewman www.cheryllewman.com

from pptx import Presentation
from pptx.util import Inches, Pt    
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# global variables
imagesList = ['ocean', 'forest', 'waterfall']
keywords = []
testImage = False
inputTitle = ''
inputText = ''
textImages = ''
titleImages = ''
prs = Presentation()

# titles/paragraphs
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
tleft = ttop = twidth = theight = Inches(3.5)
txBox = slide.shapes.add_textbox(tleft, ttop, twidth, theight)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(18)


def imagePrint(testImage):
    newlist = []
    finalImage = []
    for item in imagesList:
        for item1 in keywords:
            if item == item1:
                newlist.append(item)
                testImage = ' '.join(newlist)
                testImage += '.jpg'
                return(testImage)


createNewSlides = True
while createNewSlides:
    # title  
    inputTitle = input("What is the title of this slide? >   ")
    if inputTitle != "":
        inputTitleLC = inputTitle.lower()
        keywords = inputTitleLC.split()
        titleImages = imagePrint(testImage)
        if titleImages is None:
            titleImages = 'No Image Selected for Title.jpg'
    else:
        inputTitle = "No Image"
    # text   
    inputText = input("What text would you like on this slide? >    ")
    if inputText != "":     
        inputTextLC = inputText.lower()
        keywords = inputTextLC.split()
        textImages = imagePrint(testImage)
        if textImages is None:
            textImages = 'No Image Selected for Text.jpg'
    else:
        inputText = "No Image"
    # image
    inputImage = input("Would you like to add images to your slide? (yes/no) >    ").lower()
    if inputImage == "yes":
        inputImage = True
        userImageInput = "yes"
    else:
        break
            
    while inputImage == True:
        nl = '\n'
        # to select images for slide
        if userImageInput == "yes":
            if titleImages != '' or textImages != '':
                userChoice = input(f"Choose the image to add: {nl} {titleImages.strip('.jpg')} {nl} {textImages.strip('.jpg')} {nl} --------------- {nl} {nl} ").lower()
                if userChoice == titleImages.strip('.jpg'):
                    img_path = titleImages
                    blank_slide_layout = prs.slide_layouts[0]
                    top = Inches(1)
                    left = Inches(1)
                    height = Inches(1.5)
                    pic = slide.shapes.add_picture(img_path, left, top, height=height)  
                    titleImages = ''   
                    print('You have added the image.')              
                elif userChoice == textImages.strip('.jpg'):
                    img_path = textImages
                    blank_slide_layout = prs.slide_layouts[0]
                    left = Inches(6)
                    top = Inches(5)
                    height = Inches(1.5)
                    pic2 = slide.shapes.add_picture(img_path, left, top, height=height)  
                    textImages = ''
                    print('You have added the image.')              
            
            else:
                print('You have added all possible images for this slide.')
                break

        # to exit loop
        inputImage = input("Would you like to insert another image? (yes/no) >    ")
        if inputImage == "yes":
            inputImage = True
        else:
            break
    break

# build slide 
createSlide = print("Creating your slide now!")

title.text = inputTitle
p.text = inputText

prs.save('presentation.pptx')
